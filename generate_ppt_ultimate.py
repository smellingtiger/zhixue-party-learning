# -*- coding: utf-8 -*-
"""
智学党建学习平台 - 展会宣传PPT生成器 (终极版 - 多屏堆叠展示)
第3页建设价值 + 第4页整体架构 增加大量产品截图矩阵展示
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import glob


def set_slide_background(slide, r, g, b):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_text_box(slide, text, left, top, width, height, font_size=18,
                 bold=False, color=(255, 255, 255), alignment=PP_ALIGN.LEFT,
                 font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_image_safe(slide, image_path, left, top, width, height):
    """安全添加图片"""
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width, height)
            return True
        except Exception as e:
            print(f"    [WARN] 图片添加失败: {e}")
            return False
    return False


def find_screenshot(pattern, screenshot_dir):
    """查找匹配的截图文件"""
    matches = glob.glob(os.path.join(screenshot_dir, pattern))
    return matches[0] if matches else None


def create_cover_page(prs, screenshot_dir):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    add_text_box(slide, "数据驱动  精准赋能",
                 Inches(0.8), Inches(1.8), Inches(11.73), Inches(1.0),
                 font_size=52, bold=True, color=(0, 200, 255), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, "校本大模型全栈解决方案",
                 Inches(0.8), Inches(2.9), Inches(11.73), Inches(0.7),
                 font_size=34, color=(255, 255, 255), alignment=PP_ALIGN.CENTER)

    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.7), Inches(4.33), Inches(0.03))
    divider.fill.solid()
    divider.line.fill.background()

    add_text_box(slide, "智学党建学习平台",
                 Inches(0.8), Inches(4.0), Inches(11.73), Inches(0.5),
                 font_size=26, color=(100, 180, 255), alignment=PP_ALIGN.CENTER)

    # 首页截图
    img = find_screenshot("homepage*.png", screenshot_dir)
    if img:
        add_image_safe(slide, img, Inches(2.5), Inches(4.7), Inches(8.33), Inches(2.5))

    add_text_box(slide, "AI + Education | Data-Driven | Intelligent Empowerment",
                 Inches(0.8), Inches(6.9), Inches(11.73), Inches(0.35),
                 font_size=13, color=(150, 150, 150), alignment=PP_ALIGN.CENTER)
    return slide


def create_toc_page(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    add_text_box(slide, "目录", Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
                 font_size=42, bold=True, color=(0, 200, 255))

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.line.fill.background()

    toc_items = [
        ("01", "政策导向", "人工智能+教育已写入国家顶层设计"),
        ("02", "AI应用需求", "教师负担重、个性化教学难、AI大模型是关键钥匙"),
        ("03", "AI精准应用", "GenAI技术从炫技走向专业，校本垂域模型是基础"),
    ]

    y_pos = 1.8
    for num, title, desc in toc_items:
        num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(0.8), Inches(0.8))
        num_box.fill.solid()
        add_text_box(slide, title, Inches(1.9), Inches(y_pos + 0.05), Inches(3), Inches(0.5),
                     font_size=26, bold=True, color=(0, 220, 255))
        add_text_box(slide, desc, Inches(1.9), Inches(y_pos + 0.5), Inches(4.5), Inches(0.5),
                     font_size=14, color=(200, 200, 200))
        y_pos += 1.5
    return slide


def create_solution_overview_enhanced(prs, screenshot_dir):
    """
    第3页 - 建设价值（增强版）
    上半部分：流程步骤 + 右侧"建设价值"大字
    下半部分：6张产品截图矩阵 (2行 x 3列) 展示实际产品界面
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    # 标题
    add_text_box(slide, "建设价值", Inches(0.5), Inches(0.25), Inches(5), Inches(0.55),
                 font_size=34, bold=True, color=(0, 200, 255))

    # 流程步骤（上半部分 - 紧凑版）
    steps = ["教学数据\n与资源", "自动语料\n提取", "数据集\n构建", "垂直模型\n微调",
             "AI场景化\n应用", "新数据\n回流", "模型持续\n进化"]

    x_start = 0.35
    step_width = 1.75
    y_step = 0.95

    for i, text in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x_start + i * step_width), Inches(y_step),
                                      Inches(1.6), Inches(0.85))
        box.fill.solid()
        add_text_box(slide, text, Inches(x_start + i * step_width + 0.05), Inches(y_step + 0.15),
                    Inches(1.5), Inches(0.65), font_size=11, color=(255, 255, 255), alignment=PP_ALIGN.CENTER)

    # 右侧"建设价值"竖排大字
    value_chars = ["建", "设", "价", "值"]
    y_val = 0.85
    for char in value_chars:
        add_text_box(slide, char, Inches(12.0), Inches(y_val), Inches(1.0), Inches(0.75),
                     font_size=36, bold=True, color=(0, 180, 255), alignment=PP_ALIGN.CENTER)
        y_val += 0.8

    # ===== 下半部分：6张产品截图矩阵 (2行 x 3列) =====
    matrix_y_start = 2.0
    img_w = 4.05
    img_h = 2.55
    gap_x = 0.18
    gap_y = 0.18

    # 截图配置：(文件名模式, 标签)
    screenshots_config = [
        # 第一行：核心平台页面
        ("home_page*.png", "智能体首页"),
        ("library_page*.png", "资源汇聚"),
        ("knowledge_base*.png", "知识库管理"),
        # 第二行：智能体应用页面
        ("ai_course*.png", "AI助教-课程"),
        ("ai_profile*.png", "AI助学-画像"),
        ("training_candidates*.png", "AI助评-培训"),
    ]

    for idx, (pattern, label) in enumerate(screenshots_config):
        row = idx // 3
        col = idx % 3
        x = 0.45 + col * (img_w + gap_x)
        y = matrix_y_start + row * (img_h + gap_y)

        # 查找并添加截图
        img_path = find_screenshot(pattern, screenshot_dir)
        if img_path:
            add_image_safe(slide, img_path, Inches(x), Inches(y), Inches(img_w), Inches(img_h))
        else:
            # 占位符
            placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(img_w), Inches(img_h))
            placeholder.fill.solid()
            add_text_box(slide, f"[{label}]", Inches(x + 0.3), Inches(y + img_h/2 - 0.25),
                        Inches(img_w - 0.6), Inches(0.5), font_size=16, color=(180, 180, 180), alignment=PP_ALIGN.CENTER)

        # 底部标签条
        label_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(x), Inches(y + img_h - 0.02), Inches(img_w), Inches(0.32))
        label_bar.fill.solid()
        add_text_box(slide, label, Inches(x + 0.1), Inches(y + img_h + 0.03),
                    Inches(img_w - 0.2), Inches(0.28), font_size=11, bold=True,
                    color=(0, 220, 255), alignment=PP_ALIGN.CENTER)

    return slide


def create_architecture_enhanced(prs, screenshot_dir):
    """
    第4页 - 整体架构（增强版）
    三层架构，每层配有对应的产品截图展示
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    # 标题
    add_text_box(slide, "整体架构 - 产品全景", Inches(0.5), Inches(0.2), Inches(7), Inches(0.55),
                 font_size=32, bold=True, color=(0, 200, 255))

    # 三层架构配置
    layers = [
        {
            "title": "L1 基于垂域模型驱动的AI应用智能体",
            "color": (0, 120, 215),
            "y": 0.85,
            "screenshots": [
                ("ai_course*.png", "AI助教"),
                ("ai_profile*.png", "AI助学"),
                ("training_candidates*.png", "AI助评"),
            ]
        },
        {
            "title": "L2 AI垂域训练平台 | 构建数据集 · 训练垂域模型",
            "color": (0, 80, 180),
            "y": 2.95,
            "screenshots": [
                ("knowledge_base*.png", "知识库构建"),
                ("course_learn_detail*.png", "模型训练"),
                ("notes_page*.png", "数据资产标注"),
            ]
        },
        {
            "title": "L3 AI多模态资源中心 | 汇聚 · 分析 · 挖掘校内优质数据资产",
            "color": (0, 50, 140),
            "y": 5.05,
            "screenshots": [
                ("library_page*.png", "资源汇聚中心"),
                ("bookshelf_page*.png", "书架资源管理"),
                ("profile_page*.png", "个人数据画像"),
            ]
        },
    ]

    layer_height = 1.9
    label_width = 3.8
    img_area_start_x = 4.1
    img_w = 2.95
    img_h = 1.65
    img_gap = 0.12

    for layer in layers:
        y = layer["y"]
        color = layer["color"]

        # 左侧：层级标签
        label_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(0.35), Inches(y), Inches(label_width), Inches(layer_height))
        label_box.fill.solid()

        # 层级标题
        add_text_box(slide, layer["title"], Inches(0.5), Inches(y + 0.25), Inches(label_width - 0.3), Inches(0.55),
                     font_size=15, bold=True, color=(255, 255, 255), alignment=PP_ALIGN.CENTER)

        # 层级描述（根据层级不同显示不同内容）
        if "L1" in layer["title"]:
            desc = "三大智能体:\n• 助教智能体\n• 助学智能体\n• 助评智能体"
        elif "L2" in layer["title"]:
            desc = "核心能力:\n• 数据集构建\n• 模型微调\n• 智能体编排"
        else:
            desc = "数据资产:\n• 多模态资源\n• 知识图谱\n• 校本数据"

        add_text_box(slide, desc, Inches(0.5), Inches(y + 0.85), Inches(label_width - 0.3), Inches(0.95),
                     font_size=11, color=(210, 210, 210), alignment=PP_ALIGN.LEFT)

        # 右侧：该层对应的产品截图（3张横排）
        for idx, (pattern, label) in enumerate(layer["screenshots"]):
            x = img_area_start_x + idx * (img_w + img_gap)
            img_path = find_screenshot(pattern, screenshot_dir)

            if img_path:
                add_image_safe(slide, img_path, Inches(x), Inches(y + 0.12), Inches(img_w), Inches(img_h))
            else:
                ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.12), Inches(img_w), Inches(img_h))
                ph.fill.solid()
                add_text_box(slide, f"[{label}]", Inches(x + 0.2), Inches(y + 0.12 + img_h/2 - 0.2),
                            Inches(img_w - 0.4), Inches(0.4), font_size=13, color=(170, 170, 170), alignment=PP_ALIGN.CENTER)

            # 截图底部小标签
            tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(x), Inches(y + layer_height - 0.32), Inches(img_w), Inches(0.28))
            tag.fill.solid()
            add_text_box(slide, label, Inches(x + 0.05), Inches(y + layer_height - 0.28),
                        Inches(img_w - 0.1), Inches(0.24), font_size=10, bold=True,
                        color=(0, 220, 255), alignment=PP_ALIGN.CENTER)

    # 右侧循环箭头示意
    add_text_box(slide, "↻\n持\n续\n优\n化", Inches(12.5), Inches(2.8), Inches(0.7), Inches(2.2),
                 font_size=14, color=(0, 200, 255), alignment=PP_ALIGN.CENTER)

    return slide


def create_product_01(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    add_text_box(slide, "方案产品 — 多模态资源中心",
                 Inches(0.8), Inches(1.5), Inches(8), Inches(1),
                 font_size=44, bold=True, color=(0, 220, 255))
    add_text_box(slide, "01", Inches(9), Inches(1.5), Inches(2.5), Inches(2),
                 font_size=96, bold=True, color=(0, 150, 220))
    add_text_box(slide, "本地化部署，资源智能融合让校本数据从沉睡到觉醒！",
                 Inches(0.8), Inches(2.8), Inches(8), Inches(0.6),
                 font_size=20, color=(200, 200, 200))
    add_text_box(slide, "校本筑基 · 数智赋能 · 全链贯通",
                 Inches(8), Inches(6.2), Inches(4.5), Inches(0.4),
                 font_size=14, color=(100, 180, 255), alignment=PP_ALIGN.RIGHT)
    return slide


def create_product_01_detail(prs, screenshot_dir):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    add_text_box(slide, "数字资源汇聚与应用 - 产品实景",
                 Inches(0.5), Inches(0.25), Inches(7), Inches(0.55),
                 font_size=30, bold=True, color=(0, 200, 255))

    modules = [
        ("数据采集&存储", ["公共数据资源", "部门数据资源", "群组数据资源", "个人数据资源"]),
        ("数据治理", ["数据清洗", "多语种转写", "人像识别", "多模态分析"]),
        ("数据服务&应用", ["智能云搜", "语音检索", "OCR检索", "知识图谱"]),
        ("数据智能", ["知识图谱", "大模型", "数字人", "AI智能"])
    ]

    x_positions = [0.3, 3.4, 6.5, 9.6]

    for idx, (title, items) in enumerate(modules):
        x = x_positions[idx]
        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(0.95), Inches(2.9), Inches(0.48))
        header.fill.solid()
        add_text_box(slide, title, Inches(x + 0.1), Inches(1.02), Inches(2.7), Inches(0.38),
                     font_size=13, bold=True, color=(255, 255, 255), alignment=PP_ALIGN.CENTER)

        y_item = 1.58
        for item in items:
            item_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(y_item), Inches(2.6), Inches(0.38))
            item_box.fill.solid()
            add_text_box(slide, item, Inches(x + 0.25), Inches(y_item + 0.06), Inches(2.4), Inches(0.28),
                         font_size=11, color=(255, 255, 255), alignment=PP_ALIGN.CENTER)
            y_item += 0.46

    # Library截图
    img = find_screenshot("library_page*.png", screenshot_dir)
    if img:
        add_image_safe(slide, img, Inches(0.5), Inches(3.75), Inches(12.33), Inches(3.55))
    else:
        add_text_box(slide, "[此处显示 图书馆/资源中心 产品截图]",
                    Inches(0.5), Inches(4.8), Inches(12.33), Inches(1.5),
                    font_size=18, color=(150, 150, 150), alignment=PP_ALIGN.CENTER)
    return slide


def create_product_02(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    add_text_box(slide, "方案产品 — 垂域智能体创作平台",
                 Inches(0.8), Inches(1.5), Inches(8), Inches(1),
                 font_size=44, bold=True, color=(0, 220, 255))
    add_text_box(slide, "02", Inches(9), Inches(1.5), Inches(2.5), Inches(2),
                 font_size=96, bold=True, color=(0, 150, 220))
    add_text_box(slide, "从数据到应用，一站式垂域模型构建，让教育AI开发零门槛！",
                 Inches(0.8), Inches(2.8), Inches(8), Inches(0.6),
                 font_size=20, color=(200, 200, 200))
    add_text_box(slide, "校本筑基 · 数智赋能 · 全链贯通",
                 Inches(8), Inches(6.2), Inches(4.5), Inches(0.4),
                 font_size=14, color=(100, 180, 255), alignment=PP_ALIGN.RIGHT)
    return slide


def create_product_02_detail(prs, screenshot_dir):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    add_text_box(slide, "快速构建数据集 · 知识库 - 产品实景",
                 Inches(0.5), Inches(0.25), Inches(7), Inches(0.55),
                 font_size=30, bold=True, color=(0, 200, 255))

    features = [
        ("数据集构建", "支持多种格式导入\n自动标注与清洗\n版本管理"),
        ("知识库管理", "文档解析入库\n智能分片向量化\n语义检索"),
        ("模型微调", "可视化参数配置\nLoRA高效微调\n效果评估对比"),
        ("智能体编排", "拖拽式工作流\n多工具集成\n一键发布部署")
    ]

    positions = [(0.5, 0.95), (3.5, 0.95), (6.5, 0.95), (9.5, 0.95)]

    for idx, ((title, desc), (x, y)) in enumerate(zip(features, positions)):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(2.5))
        card.fill.solid()
        add_text_box(slide, title, Inches(x + 0.2), Inches(y + 0.18), Inches(2.4), Inches(0.42),
                     font_size=16, bold=True, color=(0, 220, 255), alignment=PP_ALIGN.CENTER)
        add_text_box(slide, desc, Inches(x + 0.2), Inches(y + 0.68), Inches(2.4), Inches(1.7),
                     font_size=12, color=(220, 220, 220), alignment=PP_ALIGN.CENTER)

    img = find_screenshot("knowledge_base*.png", screenshot_dir)
    if img:
        add_image_safe(slide, img, Inches(0.5), Inches(3.7), Inches(12.33), Inches(3.6))
    else:
        add_text_box(slide, "[此处显示 知识库管理 产品截图]",
                    Inches(0.5), Inches(4.8), Inches(12.33), Inches(1.5),
                    font_size=18, color=(150, 150, 150), alignment=PP_ALIGN.CENTER)
    return slide


def create_product_03(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    add_text_box(slide, "方案产品 — 智能体应用",
                 Inches(0.8), Inches(1.5), Inches(8), Inches(1),
                 font_size=44, bold=True, color=(0, 220, 255))
    add_text_box(slide, "03", Inches(9), Inches(1.5), Inches(2.5), Inches(2),
                 font_size=96, bold=True, color=(0, 150, 220))
    add_text_box(slide, "从千人一面到千人千面——让AI成为师生的超级助手！",
                 Inches(0.8), Inches(2.8), Inches(8), Inches(0.6),
                 font_size=20, color=(200, 200, 200))
    add_text_box(slide, "校本筑基 · 数智赋能 · 全链贯通",
                 Inches(8), Inches(6.2), Inches(4.5), Inches(0.4),
                 font_size=14, color=(100, 180, 255), alignment=PP_ALIGN.RIGHT)
    return slide


def create_application_scenarios(prs, screenshot_dir):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    add_text_box(slide, "智能体广场 - 三大智能体应用展示",
                 Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 font_size=28, bold=True, color=(0, 200, 255))

    scenarios = [
        ("助教智能体", ["共筑AI课堂互动", "智能备课", "教学设计", "课件生成"]),
        ("助学智能体", ["个性化学习路径", "智能答疑", "学习诊断", "知识图谱导航"]),
        ("助评智能体", ["过程性评价", "智能批改", "学情分析", "成长档案"])
    ]

    x_positions = [0.3, 4.5, 8.7]

    for idx, (category, items) in enumerate(scenarios):
        x = x_positions[idx]
        cat_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(0.82), Inches(4.0), Inches(0.48))
        cat_header.fill.solid()
        add_text_box(slide, category, Inches(x + 0.1), Inches(0.9), Inches(3.8), Inches(0.38),
                     font_size=17, bold=True, color=(255, 255, 255), alignment=PP_ALIGN.CENTER)

        y_item = 1.45
        for item in items:
            item_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(y_item), Inches(3.7), Inches(0.45))
            item_box.fill.solid()
            add_text_box(slide, f"● {item}", Inches(x + 0.3), Inches(y_item + 0.08), Inches(3.4), Inches(0.32),
                         font_size=12, color=(240, 240, 240))
            y_item += 0.53

    # 底部三张智能体截图
    img_configs = [
        ("ai_course*.png", "AI助教 - 课程学习"),
        ("ai_profile*.png", "AI助学 - 个人画像"),
        ("training_candidates*.png", "AI助评 - 培训评价"),
    ]

    img_w = 4.08
    img_h = 3.15
    start_x = 0.42
    img_y = 3.85

    for i, (pattern, label) in enumerate(img_configs):
        x = start_x + i * (img_w + 0.18)
        img = find_screenshot(pattern, screenshot_dir)
        if img:
            add_image_safe(slide, img, Inches(x), Inches(img_y), Inches(img_w), Inches(img_h))
        else:
            ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(img_y), Inches(img_w), Inches(img_h))
            ph.fill.solid()
            add_text_box(slide, f"[{label}]", Inches(x + 0.3), Inches(img_y + 1.3),
                        Inches(img_w - 0.6), Inches(0.5), font_size=15, color=(170, 170, 170), alignment=PP_ALIGN.CENTER)

        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(img_y + img_h - 0.02), Inches(img_w), Inches(0.3))
        tag.fill.solid()
        add_text_box(slide, label.split(" - ")[1], Inches(x + 0.05), Inches(img_y + img_h + 0.02),
                    Inches(img_w - 0.1), Inches(0.26), font_size=11, bold=True,
                    color=(0, 220, 255), alignment=PP_ALIGN.CENTER)

    return slide


def create_advantages_page(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    add_text_box(slide, "核心优势", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6),
                 font_size=36, bold=True, color=(0, 200, 255))

    advantages = [
        ("数据安全", "本地化部署，数据不出校\n满足教育行业合规要求"),
        ("垂域专业", "针对教育场景深度优化\n模型更懂教学需求"),
        ("零代码开发", "可视化操作界面\n无需编程即可创建智能体"),
        ("持续进化", "数据回流机制\n模型能力不断优化提升")
    ]

    positions = [(0.5, 1.3), (6.7, 1.3), (0.5, 3.8), (6.7, 3.8)]

    for idx, ((title, desc), (x, y)) in enumerate(zip(advantages, positions)):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.9), Inches(2.2))
        card.fill.solid()
        add_text_box(slide, title, Inches(x + 0.3), Inches(y + 0.3), Inches(5.3), Inches(0.6),
                     font_size=22, bold=True, color=(0, 220, 255))
        add_text_box(slide, desc, Inches(x + 0.3), Inches(y + 1.0), Inches(5.3), Inches(1.0),
                     font_size=16, color=(210, 210, 210))
    return slide


def create_closing_page(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, 10, 25, 66)

    add_text_box(slide, "感谢聆听", Inches(0.8), Inches(2.5), Inches(11.73), Inches(1.2),
                 font_size=60, bold=True, color=(0, 200, 255), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, "THANK YOU FOR YOUR ATTENTION", Inches(0.8), Inches(3.8), Inches(11.73), Inches(0.6),
                 font_size=24, color=(150, 150, 150), alignment=PP_ALIGN.CENTER)

    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.6), Inches(4.33), Inches(0.03))
    divider.fill.solid()
    divider.line.fill.background()

    add_text_box(slide, "数据驱动 · 精准赋能 · 智慧教育",
                 Inches(0.8), Inches(5.0), Inches(11.73), Inches(0.5),
                 font_size=22, color=(100, 180, 255), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, "智学党建学习平台 | 校本大模型全栈解决方案",
                 Inches(0.8), Inches(5.6), Inches(11.73), Inches(0.4),
                 font_size=16, color=(120, 120, 120), alignment=PP_ALIGN.CENTER)
    return slide


def main():
    print("=" * 70)
    print("  智学党建学习平台 - 展会宣传PPT生成器 (终极版 - 多屏堆叠)")
    print("=" * 70)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    downloads_dir = os.path.join(os.path.expanduser("~"), "Downloads")
    screenshot_dir = downloads_dir

    print(f"\n[INFO] 截图目录: {screenshot_dir}")
    print("[INFO] 扫描可用截图...")

    all_patterns = [
        "homepage*.png", "home_page*.png", "library_page*.png",
        "knowledge_base*.png", "ai_course*.png", "ai_profile*.png",
        "training_candidates*.png", "profile_page*.png", "notes_page*.png",
        "bookshelf_page*.png", "course_learn_detail*.png"
    ]

    found_count = 0
    for pattern in all_patterns:
        matches = glob.glob(os.path.join(screenshot_dir, pattern))
        if matches:
            print(f"  [OK] {pattern:30s} -> {os.path.basename(matches[0])}")
            found_count += 1
        else:
            print(f"  [--] {pattern:30s} (未找到)")

    print(f"\n[INFO] 共找到 {found_count}/{len(all_patterns)} 张截图")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("\n[INFO] 开始生成幻灯片...")

    pages = [
        ("封面页（含首页截图）", lambda p: create_cover_page(p, screenshot_dir)),
        ("目录页", create_toc_page),
        ("★ 建设价值（6屏矩阵展示）", lambda p: create_solution_overview_enhanced(p, screenshot_dir)),
        ("★ 整体架构（三层+9屏展示）", lambda p: create_architecture_enhanced(p, screenshot_dir)),
        ("产品01-多模态资源中心", create_product_01),
        ("产品01详情（含Library截图）", lambda p: create_product_01_detail(p, screenshot_dir)),
        ("产品02-垂域智能体创作平台", create_product_02),
        ("产品02详情（含Knowledge Base截图）", lambda p: create_product_02_detail(p, screenshot_dir)),
        ("产品03-智能体应用", create_product_03),
        ("应用场景（三大智能体截图）", lambda p: create_application_scenarios(p, screenshot_dir)),
        ("核心优势", create_advantages_page),
        ("结束页", create_closing_page),
    ]

    for name, func in pages:
        func(prs)
        marker = "★" if "★" in name else "  "
        print(f"  {marker} {name}")

    output_filename = "智学党建学习平台_展会宣传_含截图.pptx"
    output_path = os.path.join(script_dir, output_filename)
    prs.save(output_path)

    print("\n" + "=" * 70)
    print(f"[SUCCESS] PPT生成完成!")
    print(f"[FILE]   {output_path}")
    print(f"[PAGES]  {len(pages)} 页")
    print("=" * 70)

    print("\n[DETAIL] 第3页 '建设价值' 包含:")
    print("         - 上方: 7步流程图 + '建设价值'竖排字")
    print("         - 下方: 2x3 产品截图矩阵 (6张)")
    print("           ├── 智能体首页(宣传页) / 资源汇聚 / 知识库管理")
    print("           └── AI助教-课程 / AI助学-画像 / AI助评-培训")

    print("\n[DETAIL] 第4页 '整体架构' 包含:")
    print("         - L1 应用智能体层: 3张截图 (AI助教/助学/助评)")
    print("         - L2 训练平台层:   3张截图 (知识库/训练/数据资产标注)")
    print("         - L3 资源中心层:   3张截图 (资源汇聚/书架管理/个人画像)")

    print("\n[STYLE] 深蓝科技风 (#0A1942) + 青色主题 (#00C8FF) + 16:9宽屏")
    print("\n可直接用于展会展示!")

    return output_path


if __name__ == "__main__":
    main()
